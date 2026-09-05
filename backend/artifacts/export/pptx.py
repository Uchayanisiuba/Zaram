"""HTML → .pptx.

`CLAUDE.md` deferred PowerPoint with "later, only if asked". It was asked for on
10 August 2026, so it is here — and the way it is built is the interesting part.

**Headings are the slide boundaries.** There is no separate deck format and no
second authoring path: `<h1>` becomes the title slide, every `<h2>` starts a new
one, and the paragraphs and list items beneath it become that slide's bullets.
That falls out of HTML being the source of truth rather than being designed on
top of it, and it buys something a dedicated deck format would not — **any
document Zaram has already generated can be exported as slides**, because every
document has headings. A proposal becomes a pitch without being rewritten.

What that costs, stated plainly
-------------------------------
`_reader` collects tables separately from blocks, so their position relative to
the headings is not recoverable. Tables therefore land on their own slides at
the end rather than inside the section they belonged to. That is a real loss of
ordering and it is preferred to the alternatives: dropping them silently, or
teaching the reader to interleave for one exporter's benefit. A deck is
re-ordered by hand anyway, and a missing table is not.

It carries the same design as everything else — 4 September 2026
--------------------------------------------------------------
This used to read *"nothing here is a theme… a deck someone presents is one they
will restyle, and a hardcoded palette is one more thing to undo."* The
maintainer looked at the output and disagreed, and the measurement settles it
rather than the argument: the deck came out **4:3** — 10 by 7.5 inches, the
default `python-pptx` template — with **not one run carrying a font, a size or a
colour**. That is not restraint leaving room for someone's own template; it is
Office 2003 with the lights off.

The reasoning was wrong in a specific way worth naming. *Not styling* is not the
neutral choice — python-pptx's stock template is itself a design, just nobody's.
`theme.py` was written after exactly this discovery about Word, where the same
"leave it to the renderer" produced Calibri 11 and Word 2007 blue, and its
opening line applies here unchanged: **two renderers had two designs, and only
one of them was designed.** So the deck reads the same tokens as the HTML and
the Word file, and a proposal exported three ways looks like three views of one
document.

**16:9, because 4:3 is not a taste question.** Every projector, laptop and
conferencing tool made since about 2012 is widescreen; a 4:3 deck is pillarboxed
on all of them, which reads as a file made by accident.

What is still left alone: the layouts themselves. Slide masters, backgrounds and
placeholder geometry stay the template's, so a user who applies their own theme
in PowerPoint gets it applied to a deck whose structure is ordinary. Type and
colour are what carry the identity; moving the boxes would be the "one more
thing to undo" the old note was actually right about.
"""

from __future__ import annotations

import io
from typing import List, Tuple

from . import _reader
from .. import theme
from .base import Availability, module_available

#: Layout indices in python-pptx's default template. Named because `slide_layouts[1]`
#: at a call site is unreadable and wrong-by-one is silent.
_TITLE_SLIDE = 0
_TITLE_AND_CONTENT = 1
_TITLE_ONLY = 5

#: 16:9, in inches. The height is the template's own 7.5; only the width moves,
#: so every placeholder the layouts define keeps its vertical position.
_WIDESCREEN_IN = (13.333, 7.5)

#: Type sizes for a room, not for a page.
#:
#: `theme.py`'s sizes are print sizes — 11pt body, 20pt title — and a deck read
#: from the back of a room is a different instrument. What is shared is the
#: *faces* and the *colours*, which is what makes the three exports look like
#: one document; the scale is this format's own.
_TITLE_PT = 32.0
_COVER_TITLE_PT = 40.0
_BULLET_PT = 18.0
_TABLE_PT = 12.0

#: Past this, a slide is a wall of text nobody reads from the back of a room.
#: Overflow continues on a slide marked "(cont.)" rather than being dropped or
#: shrunk to fit — losing content silently is the one outcome that is not
#: recoverable by the person editing the deck afterwards.
_MAX_BULLETS = 8


class PptxExporter:
    extension = "pptx"
    media_type = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    label = "PowerPoint"

    def availability(self) -> Availability:
        return module_available("pptx", needed_for="PowerPoint export")

    def export(self, document_html: str, *, filename: str = "") -> bytes:
        from pptx import Presentation

        from pptx.util import Inches

        doc = _reader.read(document_html)
        deck = Presentation()
        deck.slide_width = Inches(_WIDESCREEN_IN[0])
        deck.slide_height = Inches(_WIDESCREEN_IN[1])

        title, slides = self._outline(doc)

        cover = deck.slides.add_slide(deck.slide_layouts[_TITLE_SLIDE])
        cover.shapes.title.text = title or "Untitled"
        self._style(cover.shapes.title, size=_COVER_TITLE_PT, colour=theme.INK, bold=True)
        # Placeholder 1 is the subtitle. Cleared rather than left holding
        # "Click to add subtitle", which is what ships in the template and what
        # would otherwise be presented to a room.
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = ""

        for heading, bullets in slides:
            for index, chunk in enumerate(self._chunk(bullets)):
                slide = deck.slides.add_slide(deck.slide_layouts[_TITLE_AND_CONTENT])
                slide.shapes.title.text = heading if index == 0 else f"{heading} (cont.)"
                # The accent on section titles and nowhere else, which is the
                # same rule `theme.ACCENT` follows on the page: one accent, used
                # for the thing that says what this is.
                self._style(slide.shapes.title, size=_TITLE_PT, colour=theme.ACCENT, bold=True)
                frame = slide.placeholders[1].text_frame
                frame.clear()
                for position, bullet in enumerate(chunk):
                    paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
                self._style(slide.placeholders[1], size=_BULLET_PT, colour=theme.INK)

        for table in doc.tables:
            self._table_slide(deck, table)

        buffer = io.BytesIO()
        deck.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def _outline(doc: _reader.Document) -> Tuple[str, List[Tuple[str, List[str]]]]:
        """Title, then (heading, bullets) per slide.

        Content before the first `<h2>` is gathered under the document title, so
        an opening paragraph is not thrown away for arriving early.
        """
        title = ""
        slides: List[Tuple[str, List[str]]] = []
        current: Tuple[str, List[str]] | None = None

        for block in doc.body_blocks():
            text = block.text.strip()
            if not text:
                continue

            if block.tag == "h1" and not title:
                title = text
                continue

            if block.tag in ("h1", "h2", "h3"):
                current = (text, [])
                slides.append(current)
                continue

            if current is None:
                current = (title or "Overview", [])
                slides.append(current)
            current[1].append(text)

        # A heading with nothing under it is still a slide: it is a section
        # marker, and dropping it loses the deck's structure.
        return title, slides

    @staticmethod
    def _style(shape, *, size: float, colour: str, bold: bool = False) -> None:
        """Set face, size and colour on every run in a shape.

        **Every run, and every paragraph, including the empty ones.** A run is
        created when text is assigned, so a paragraph styled before its text
        exists keeps the template's Calibri — which is how a deck ends up
        styled everywhere except the one line somebody actually reads. Setting
        it on the paragraph's own font as well covers a paragraph that has no
        runs yet.

        `theme.WORD_SERIF` rather than the CSS stack, for the reason `theme.py`
        gives about Word: PowerPoint stores one name and substitutes silently
        if it is absent, so a stack is meaningless and the honest choice is the
        face most likely to be on the machine that opens the file.
        """
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        frame = getattr(shape, "text_frame", None)
        if frame is None:
            return
        rgb = RGBColor.from_string(colour.upper())
        for paragraph in frame.paragraphs:
            for font in [paragraph.font, *(run.font for run in paragraph.runs)]:
                font.name = theme.WORD_SERIF
                font.size = Pt(size)
                font.bold = bold
                font.color.rgb = rgb

    @staticmethod
    def _chunk(bullets: List[str]) -> List[List[str]]:
        if not bullets:
            return [[]]
        return [bullets[i : i + _MAX_BULLETS] for i in range(0, len(bullets), _MAX_BULLETS)]

    @staticmethod
    def _table_slide(deck, table: _reader.Table) -> None:
        from pptx.util import Inches

        grid = ([table.header] if table.header else []) + [list(r) for r in table.rows]
        if not grid:
            return

        columns = max(len(row) for row in grid)
        slide = deck.slides.add_slide(deck.slide_layouts[_TITLE_ONLY])
        slide.shapes.title.text = table.caption or "Table"
        # Styled like every other section title. Measured and missed the first
        # time: the deck came out themed on four slides and stock on the fifth,
        # which is worse than uniformly stock because it reads as a rendering
        # fault rather than as a plain template.
        PptxExporter._style(
            slide.shapes.title, size=_TITLE_PT, colour=theme.ACCENT, bold=True
        )

        shape = slide.shapes.add_table(
            len(grid), columns, Inches(0.6), Inches(1.8), Inches(9), Inches(0.4 * len(grid))
        )
        for r, row in enumerate(grid):
            for c in range(columns):
                cell = shape.table.cell(r, c)
                cell.text = str(row[c]) if c < len(row) else ""
                # A header row is set in the muted grey the page uses for the
                # same job, rather than in the template's reversed-out white on
                # a solid band. Reading a table off a slide is scanning, and the
                # page already decided what scanning looks like here.
                PptxExporter._style(
                    cell,
                    size=_TABLE_PT,
                    colour=theme.MUTED if (r == 0 and table.header) else theme.INK,
                    bold=bool(r == 0 and table.header),
                )
