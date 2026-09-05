"""HTML → slides.

The design claim worth testing is that there is **no separate deck format**:
headings are the slide boundaries, so any document Zaram has already generated
can be exported as slides without being rewritten. If that holds, a proposal
becomes a pitch for free; if it does not, `.pptx` is a second authoring path
and the pipeline rule has quietly gained an exception.
"""

from __future__ import annotations

import io

import pytest

from artifacts import export
from artifacts.html import render_deck, render_document, render_spreadsheet

pytest.importorskip("pptx", reason="python-pptx is not installed")

DECK = render_deck(
    title="Northwind — Q3",
    subtitle="Where the work stands",
    slides=[
        ("What we agreed", ["Three design days", "Revisions included"]),
        ("What is outstanding", ["Payment on 9 September"]),
        ("Next", []),
    ],
)

PROSE = render_document(
    title="A proposal",
    blocks=["Opening context before any heading."],
)

TABLE = render_spreadsheet(
    title="Q3 invoices",
    header=["Client", "Amount"],
    rows=[["Northwind", "1,470.50"]],
    caption="Outstanding",
)


def slides_of(document_html: str):
    from pptx import Presentation

    return Presentation(io.BytesIO(export.render(document_html, "pptx"))).slides


def text_of(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


class TestTheOutline:
    def test_the_title_becomes_the_first_slide(self):
        slides = slides_of(DECK)

        assert slides[0].shapes.title.text == "Northwind — Q3"

    def test_each_heading_becomes_a_slide(self):
        titles = [s.shapes.title.text for s in slides_of(DECK)]

        assert "What we agreed" in titles
        assert "What is outstanding" in titles

    def test_bullets_land_under_their_heading(self):
        slide = next(s for s in slides_of(DECK) if s.shapes.title.text == "What we agreed")

        body = text_of(slide)
        assert "Three design days" in body
        assert "Revisions included" in body
        # Not leaked onto the wrong slide.
        assert "Payment on 9 September" not in body

    def test_an_empty_heading_is_still_a_slide(self):
        # A section marker. Dropping it loses the deck's structure.
        assert "Next" in [s.shapes.title.text for s in slides_of(DECK)]

    def test_the_subtitle_placeholder_is_not_left_saying_click_to_add(self):
        """What ships in python-pptx's template, and what would be presented.

        An unfilled placeholder is not blank — it carries prompt text, and a
        deck opened in front of a room showing "Click to add subtitle" is the
        kind of detail that reads as unfinished software.
        """
        assert "Click to add" not in text_of(slides_of(DECK)[0])


class TestAnyDocumentBecomesADeck:
    def test_a_plain_document_exports_without_being_rewritten(self):
        """The claim the whole design rests on.

        No deck kind required, no second authoring path — a document has
        headings, and headings are slides.
        """
        slides = slides_of(PROSE)

        assert len(slides) >= 1
        assert slides[0].shapes.title.text == "A proposal"

    def test_content_before_the_first_heading_is_kept(self):
        # Gathered under the title rather than thrown away for arriving early.
        body = "\n".join(text_of(s) for s in slides_of(PROSE))

        assert "Opening context before any heading." in body


class TestTables:
    def test_a_table_becomes_a_real_table_not_a_paragraph(self):
        slides = slides_of(TABLE)
        tables = [sh for s in slides for sh in s.shapes if sh.has_table]

        assert tables, "the table was dropped or flattened into text"
        grid = tables[0].table
        assert grid.cell(0, 0).text == "Client"
        assert grid.cell(1, 0).text == "Northwind"
